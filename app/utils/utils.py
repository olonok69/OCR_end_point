import base64
import io
import os
from PIL import Image
from pdf2image import convert_from_bytes
import numpy as np
import docx2txt
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup as BS
import copy
from fastapi import FastAPI
from detectaicore import Job, image_file_names
import zipfile
from pathlib import Path

try:
    from image.ocr.app.utils.ocr_c import (
        pil2PIX32,
    )


except:
    from utils.ocr_c import (
        pil2PIX32,
    )


ROOT_DIR = os.path.dirname(os.path.abspath(__file__))


def ocr_image(im, app):
    """
    Extract text from an Image using Leptonica and tessetact
    params:
    im: PIL image
    app: FastAPI endpoint. Contains global objects with tesseract and Leptonica C objects
    return:
    text extracted
    """
    pix = pil2PIX32(im, app.leptonica, app.ffi)

    # Get information about DPI
    x_dpi = app.ffi.new("int *")
    y_dpi = app.ffi.new("int *")
    app.leptonica.pixGetResolution(pix, x_dpi, y_dpi)
    app.tesseract.TessBaseAPISetImage2(app.api, pix)
    app.tesseract.TessBaseAPIRecognize(app.api, app.ffi.NULL)
    # Print whole recognized text
    utf8_text = app.ffi.string(app.tesseract.TessBaseAPIGetUTF8Text(app.api)).decode(
        "utf-8"
    )
    return utf8_text


def extract_from_single_image(data, app):
    """
    Process imagenes and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """
    im_b64 = data.get("source").get("content")
    ext = data.get("source").get("file_type")

    img_bytes = base64.b64decode(im_b64.encode("utf-8"))
    # Open Image
    im = Image.open(io.BytesIO(img_bytes))
    if ext == "gif":
        mypalette = im.getpalette()
        im.putpalette(mypalette)
        new_im = Image.new("RGBA", im.size)
        new_im.paste(im)
        im = copy.deepcopy(new_im)
    # extract text
    utf8_text = ocr_image(im, app)
    data["source"]["content"] = utf8_text
    return utf8_text, data


def extract_text_from_pdf(data, app):
    """
    extract and Process images in a pdf file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    text = {}
    im_b64 = data.get("source").get("content")

    img_bytes = base64.b64decode(im_b64.encode("utf-8"))
    # Open Image
    pdf_file = convert_from_bytes(img_bytes, dpi=300, grayscale=True)
    for i, page in enumerate(pdf_file):
        page_arr = np.asarray(page)
        im = Image.fromarray(page_arr)
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    return txt[: -len(mark)]


def extract_text_from_doc(data, app):
    """
    extract and Process images in a doc/docx file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")

    # write docx to temporal folder
    temp_docx = os.path.join(directory, "temp.docx")

    with open(temp_docx, "wb") as f:
        f.write(base64.b64decode(im_b64))

    _ = docx2txt.process(temp_docx, images)
    files = os.listdir(images)

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    return txt[: -len(mark)]


def extract_text_from_ppt(data, app):
    """
    extract and Process images in a ppt/pptx file and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    def write_image(shape, images, n):
        image = shape.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = os.path.join(images, "ppt_image{:03d}.{}".format(n, image.ext))
        n += 1

        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return n

    def visitor(shape, images, n):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                n = visitor(s, images, n)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n = write_image(shape, images, n)
        return n

    def iter_picture_shapes(prs, n):
        for slide in prs.slides:
            for shape in slide.shapes:
                n = visitor(shape, images, n)

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    temp_pptx = os.path.join(directory, "temp.pptx")

    with open(temp_pptx, "wb") as f:
        f.write(base64.b64decode(im_b64))
    # Open PPT and extract Images
    iter_picture_shapes(Presentation(temp_pptx), n)
    # list all images extracted in images
    files = os.listdir(images)

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    return txt[: -len(mark)]


def extract_text_from_html(data, app):
    """
    extract and Process images in a html file with images embedded in base64 in src and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    html_temp = os.path.join(directory, "temp.html")

    with open(html_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))
    # Open html and extract Images
    with open(html_temp) as html_wr:
        html_data = html_wr.read()

    soup = BS(html_data, features="lxml")

    images_arr = soup.find_all("img")

    for ind, i in zip(images_arr, range(len(images_arr))):
        image_data_base64 = ind["src"].split(",")[1]
        decoded_img_data = base64.b64decode(image_data_base64)
        with open(os.path.join(images, f"site_{i}.png"), "wb+") as img_wr:
            img_wr.write(decoded_img_data)

    # list all images extracted in images
    files = os.listdir(images)

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    return txt[: -len(mark)]


def extract_text_from_odt(data, app):
    """
    extract and Process images in a odf file with images embedded in base64 in src and extract text via ocr
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    # page number
    n = 0
    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)

    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    odf_temp = os.path.join(directory, "test.odt")
    # write temp odf file
    with open(odf_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))
    # read the temporal odf file with zipfile
    with zipfile.ZipFile(odf_temp) as zf:
        # Get list of file names inside zip
        file_list = zf.namelist()
        # Check for Pictures folder
        pictures_folder = "Pictures"
        for file in file_list:
            if pictures_folder in file:
                # Extract Pictures folder
                zf.extract(file, path=images)

    # list all images extracted in images. Zip file look for images in Folder named Pictures and extract them under that name.
    # Odf it is a kind of zip container
    files = os.listdir(os.path.join(images, "Pictures"))

    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, "Pictures", file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    return txt[: -len(mark)]


def extract_text_from_rtf(data, app):
    """
    extract and Process images in a rtf file with images. Use of unrtf tool which has to be installed in the OS
    params:
    data: list of documents from detect AI
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint
    """

    def convert_rtf_to_doc(input_file, output_file, images, ROOT_DIR):
        os.chdir(images)
        os.system(f"/usr/bin/unrtf --html {input_file} > {output_file}")
        os.chdir(ROOT_DIR)

    # temporal directories
    directory = os.path.join(ROOT_DIR, "tmp")
    if not os.path.exists(directory):
        os.makedirs(directory)
    images = os.path.join(ROOT_DIR, "tmp", "images")
    if not os.path.exists(images):
        os.makedirs(images)

    text = {}
    im_b64 = data.get("source").get("content")
    # write docx to temporal folder
    rtf_temp = os.path.join(directory, "temp.rtf")

    with open(rtf_temp, "wb") as f:
        f.write(base64.b64decode(im_b64))

    # convert RTF to HTML
    output_file = os.path.join(images, "test_rtf.html")
    convert_rtf_to_doc(rtf_temp, output_file, images, ROOT_DIR)
    # Open html and extract Images
    with open(output_file) as html_wr:
        html_data = html_wr.read()

    soup = BS(html_data, features="lxml")
    # find all references to Images  and save Pathnames to list files
    images_arr = soup.find_all("img")
    files = []
    for ind, i in zip(images_arr, range(len(images_arr))):
        files.append(ind["src"])

    # OCR each image on the list
    for i, file in zip(range(len(files)), files):
        im = Image.open(os.path.join(images, file))
        # Extract text from Image
        utf8_text = ocr_image(im, app)
        text[i] = utf8_text
    # Final text in single file with a mark
    txt = ""
    mark = "\n### NEW PAGE ###\n"
    for key in text.keys():
        txt = txt + text[key] + mark

    # delete temp folder
    shutil.rmtree(directory, ignore_errors=True)

    return txt[: -len(mark)]


def extract_documents_from_request(document, app):
    """
    Process documents depending of their file type
    params:
    document: dictionary containing metadata and data of a document
    app global variables from FastAPI. contains tessetact and leptonica objects
    return:
    text extracted and endpoint

    """
    # get extension of document from request
    ext = document.get("source").get("file_type")

    if ext in ["jpg", "png", "jpeg", "tif", "tiff", "gif", "bmp"]:
        utf8_text, document = extract_from_single_image(document, app)
    elif ext == "pdf":
        utf8_text = extract_text_from_pdf(document, app)
    elif ext in ["doc", "docx"]:
        utf8_text = extract_text_from_doc(document, app)
    elif ext in ["ppt", "pptx"]:
        utf8_text = extract_text_from_ppt(document, app)
    elif ext in ["htm", "html"]:
        utf8_text = extract_text_from_html(document, app)
    elif ext in ["odt"]:
        utf8_text = extract_text_from_odt(document, app)
    elif ext in ["rtf"]:
        utf8_text = extract_text_from_rtf(document, app)

    return utf8_text, document


def process_request(
    list_docs: list[dict], app: FastAPI, jobs: dict, new_task: Job, cypher: int = 0
):
    """ "
    process list of base64 Image documents to OCR


    params:
    list_docs: list of documents from detect AI. List of dictionaries containing metadata and data of documents
    app global variables from FastAPI. contains tessetact and leptonica objects
    jobs: dictionary to hold status of the Job
    new_task: Job object
    cypher: boolean , True encrypt the output/False No
    return:
    processed_docs : list of dictionaries containing document processed , this is pass trought OCR and text extracted. The extracted text replace the original base64 content
    documents_non_teathred : list of dictionaries containing {id : reason of not treating this id}
    """
    jobs[new_task.uid].status = "Start Extracting Text From Documents"
    documents_non_teathred = []
    processed_docs = []
    for data in list_docs:
        file_type = data.get("source").get("file_type")
        file_name = data.get("source").get("file_name")
        print(
            f"Processing file : {file_name} length : {len(data.get('source').get('content'))}"
        )
        # if file type not valid
        if not file_type or not (file_type in image_file_names):
            documents_non_teathred.append({data.get("id"): "File type not valid"})
            print(f"File : {file_name}  No Valid Extension")
            continue

        utf8_text, data = extract_documents_from_request(data, app)
        # if no text extracted doc to documents_non_teathred
        if len(utf8_text) < 3:
            documents_non_teathred.append(
                {data.get("id"): "No Images in this document to process"}
            )
            print(f"File : {file_name}  No Text in this file")
            continue
        # create encrypted content

        if cypher:
            # chunck = encrypt(utf8_text)
            # data["crypto"] = (base64.b64encode(chunck)).decode("ascii")
            # case that we want to activate recover the file crypto.py , import it and
            print("this functionality has been deactivated")
        else:
            pass
        data["source"]["content"] = utf8_text
        processed_docs.append(data)
        jobs[
            new_task.uid
        ].status = f"Extracted Text From {len(processed_docs)} Documents"
    return processed_docs, documents_non_teathred
